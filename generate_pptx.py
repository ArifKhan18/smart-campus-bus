import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 16:9 Widescreen Presentation Colors (Exact match to HTML Classic & Dynamic palettes)
DARK_BLUE = RGBColor(26, 54, 93)        # #1a365d (deep navy blue for Title/Thank You)
BLUE = RGBColor(37, 99, 235)            # #2563eb (accent primary blue)
LIGHT_BLUE_BG = RGBColor(239, 246, 255) # #eff6ff
BLUE_BORDER = RGBColor(191, 219, 254)   # #bfdbfe

WHITE = RGBColor(255, 255, 255)         # #ffffff
DARK_TEXT = RGBColor(26, 32, 44)        # #1a202c (primary dark text)
GRAY_TEXT = RGBColor(74, 85, 104)       # #4a5568 (secondary body text)
LIGHT_GRAY = RGBColor(240, 244, 248)    # #f0f4f8 (slide background / card label strips)
BORDER_GRAY = RGBColor(226, 232, 240)   # #e2e8f0 (subtle card borders)

# Highlight Palettes
GREEN_BG = RGBColor(236, 253, 245)      # #ecfdf5
GREEN_BORDER = RGBColor(167, 243, 208)  # #a7f3d0
GREEN_TEXT = RGBColor(6, 95, 70)        # #065f46
GREEN_ACCENT = RGBColor(16, 185, 129)   # #10b981

ORANGE_BG = RGBColor(255, 251, 235)     # #fffbeb
ORANGE_BORDER = RGBColor(253, 230, 138) # #fde68a
ORANGE_TEXT = RGBColor(146, 64, 14)     # #92400e
ORANGE_ACCENT = RGBColor(245, 158, 11)  # #f59e0b

PURPLE_BG = RGBColor(245, 243, 255)     # #f5f3ff
PURPLE_BORDER = RGBColor(221, 214, 254) # #ddd6fe

IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'project picture')

def add_dark_bg(slide):
    """Add solid deep navy background for Title & Thank You slides"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

def add_light_bg(slide):
    """Add light gray background matching HTML secondary body background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY

def add_content_card(slide, left, top, width, height, bg_color=WHITE, border_color=BORDER_GRAY):
    """Add rounded rectangle card shape matching HTML CSS cards"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1)
    card.adjustments[0] = 0.03
    return card

def add_title_bar(slide, title_text, icon="", slide_num=""):
    """Add consistent top header bar matching slides.html in 16:9 widescreen"""
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(10.5), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = DARK_TEXT
    run.font.name = "Outfit"

    # Blue header accent underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.98), Inches(11.733), Emu(30000))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    # Slide number badge
    if slide_num:
        num_box = slide.shapes.add_textbox(Inches(11.3), Inches(0.38), Inches(1.233), Inches(0.5))
        tf2 = num_box.text_frame
        tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = Inches(0)
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        run2 = p2.add_run()
        run2.text = slide_num
        run2.font.size = Pt(16)
        run2.font.bold = True
        run2.font.color.rgb = GRAY_TEXT
        run2.font.name = "Inter"

def add_description(slide, text, top=Inches(1.15)):
    """Add description subtitle text below header"""
    desc_box = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.733), Inches(0.45))
    tf = desc_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.color.rgb = GRAY_TEXT
    run.font.name = "Inter"

def add_rich_bullet_points(slide, points, left, top, width, height, font_size=11, bullet_type="check"):
    """Add styled bullet points with bold titles matching HTML feature lists"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    
    for i, item in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        
        # Add Bullet
        bullet_run = p.add_run()
        if bullet_type == "check":
            bullet_run.text = "✓  "
            bullet_run.font.bold = True
            bullet_run.font.color.rgb = GREEN_ACCENT
        elif bullet_type == "num":
            bullet_run.text = f"{i+1}.  "
            bullet_run.font.bold = True
            bullet_run.font.color.rgb = ORANGE_ACCENT
        elif bullet_type == "star":
            bullet_run.text = "•  "
            bullet_run.font.bold = True
            bullet_run.font.color.rgb = BLUE
        elif bullet_type == "icon":
            bullet_run.text = f"{item[0]}  "
        
        bullet_run.font.size = Pt(font_size)
        bullet_run.font.name = "Inter"

        # Content Text (can be tuple: (bold_prefix, normal_desc) or just string)
        if isinstance(item, tuple):
            prefix = item[1] if bullet_type == "icon" else item[0]
            desc = item[2] if bullet_type == "icon" else item[1]
            
            run_bold = p.add_run()
            run_bold.text = prefix + " "
            run_bold.font.bold = True
            run_bold.font.size = Pt(font_size)
            run_bold.font.color.rgb = DARK_TEXT
            run_bold.font.name = "Inter"
            
            run_desc = p.add_run()
            run_desc.text = desc
            run_desc.font.size = Pt(font_size)
            run_desc.font.color.rgb = GRAY_TEXT
            run_desc.font.name = "Inter"
        else:
            run_desc = p.add_run()
            run_desc.text = item
            run_desc.font.size = Pt(font_size)
            run_desc.font.color.rgb = GRAY_TEXT
            run_desc.font.name = "Inter"

def add_image_safe(slide, img_name, left, top, width=None, height=None):
    """Add image with standard thin border"""
    img_path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(img_path):
        kwargs = {'image_file': img_path, 'left': left, 'top': top}
        if width: kwargs['width'] = width
        if height: kwargs['height'] = height
        pic = slide.shapes.add_picture(**kwargs)
        pic.line.color.rgb = BORDER_GRAY
        pic.line.width = Pt(1)
        return True
    else:
        print(f"[WARN] Image not found: {img_name}")
        return False

def add_screenshot_card(slide, img_name, caption, left, top, width, height, caption_height=Inches(0.42)):
    """Add screenshot card with image at top and label strip at bottom matching HTML"""
    # Outer card
    add_content_card(slide, left, top, width, height + caption_height)
    
    # Inner image
    img_top = top + Inches(0.08)
    img_left = left + Inches(0.08)
    img_width = width - Inches(0.16)
    img_height = height - Inches(0.12)
    add_image_safe(slide, img_name, img_left, img_top, width=img_width, height=img_height)
    
    # Caption strip
    cap_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + height, width, caption_height)
    cap_bg.fill.solid()
    cap_bg.fill.fore_color.rgb = LIGHT_GRAY
    cap_bg.line.fill.background()
    
    cap_box = slide.shapes.add_textbox(left, top + height + Inches(0.06), width, caption_height - Inches(0.08))
    tf = cap_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = caption.upper()
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = GRAY_TEXT
    run.font.name = "Inter"

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ==================== SLIDE 1: TITLE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)

    # Badge
    badge_box = slide.shapes.add_textbox(Inches(3.666), Inches(1.1), Inches(6.0), Inches(0.4))
    tf = badge_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "SOFTWARE DEVELOPMENT PROJECT 400"
    run.font.size = Pt(13)
    run.font.color.rgb = RGBColor(200, 220, 255)
    run.font.name = "Inter"
    run.font.bold = True

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.7), Inches(10.333), Inches(1.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Smart Campus Bus\nTracking System"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Outfit"

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.333), Inches(0.6))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "A real-time GPS-based bus tracking & campus transit management web application"
    run.font.size = Pt(15)
    run.font.color.rgb = RGBColor(200, 220, 255)
    run.font.name = "Inter"

    # Team & Supervisor Card
    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.333), Inches(2.2))
    tf = info_box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(4)
    r1 = p1.add_run()
    r1.text = "Supervisor: "
    r1.font.bold = True
    r1.font.size = Pt(14)
    r1.font.color.rgb = WHITE
    r2 = p1.add_run()
    r2.text = "Humayra Ahmed, Assistant Professor"
    r2.font.size = Pt(14)
    r2.font.color.rgb = RGBColor(220, 235, 255)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(8)
    r3 = p2.add_run()
    r3.text = "Intake 51 | Section 3"
    r3.font.bold = True
    r3.font.size = Pt(13)
    r3.font.color.rgb = RGBColor(190, 215, 250)

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r4 = p3.add_run()
    r4.text = "Md Arif Khan  •  Karnia Binte Rafique  •  Suraiya Karim  •  Prosenjit Biswas  •  Proshanta Saha"
    r4.font.size = Pt(13)
    r4.font.color.rgb = RGBColor(220, 235, 255)

    # ==================== SLIDE 2: PROBLEM & OBJECTIVES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "Problem Statement & Project Objectives", "", "02")
    add_description(slide, "Addressing campus transit inefficiencies through an automated, real-time tracking and communication ecosystem.")

    # Left: Problems
    add_content_card(slide, Inches(0.8), Inches(1.7), Inches(5.666), Inches(5.2), bg_color=ORANGE_BG, border_color=ORANGE_BORDER)
    ttl = slide.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.0), Inches(0.4))
    r = ttl.text_frame.paragraphs[0].add_run()
    r.text = "❌ Existing Challenges"
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = ORANGE_TEXT

    prob_items = [
        ("Blind Waiting:", "No live bus location or arrival ETA available at campus stops."),
        ("Unreliable Schedules:", "Paper notices with unexpected delays cause missed classes."),
        ("Communication Gap:", "No passenger-driver coordination channel during breakdown events."),
        ("Administrative Blindness:", "Zero real-time fleet overview, speed metrics, or instant safety reports.")
    ]
    add_rich_bullet_points(slide, prob_items, Inches(1.1), Inches(2.45), Inches(5.1), Inches(4.2), font_size=11, bullet_type="num")

    # Right: Objectives
    add_content_card(slide, Inches(6.866), Inches(1.7), Inches(5.666), Inches(5.2), bg_color=GREEN_BG, border_color=GREEN_BORDER)
    ttl2 = slide.shapes.add_textbox(Inches(7.166), Inches(1.9), Inches(5.0), Inches(0.4))
    r2 = ttl2.text_frame.paragraphs[0].add_run()
    r2.text = "🎯 Target Objectives & Solutions"
    r2.font.size = Pt(17)
    r2.font.bold = True
    r2.font.color.rgb = GREEN_TEXT

    obj_items = [
        ("Real-Time GPS Tracking:", "High-precision location updates broadcasted every 5s on Leaflet OSM."),
        ("Haversine ETA Engine:", "Accurate stop-wise distance and arrival time calculations in real-time."),
        ("Live Community Chat:", "Bus-specific instant communication grouped by active transit trips."),
        ("Fleet Management Portal:", "Interactive route coordinator, user role controls, and visual analytics.")
    ]
    add_rich_bullet_points(slide, obj_items, Inches(7.166), Inches(2.45), Inches(5.1), Inches(4.2), font_size=11, bullet_type="check")

    # ==================== SLIDE 3: TECH STACK & ARCHITECTURE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "System Architecture & Tech Stack", "", "03")
    add_description(slide, "A scalable, decoupled client-server architecture powered by ASP.NET Core and Firebase real-time sync.")

    techs = [
        ("🖥️ HTML5 / CSS3 / JS", "Vanilla UI, Responsive & Dynamic"),
        ("🟣 ASP.NET Core 8", "C# REST Web API & JWT Auth"),
        ("🔥 Firebase Firestore", "NoSQL Real-Time Sync SDK"),
        ("🗺️ Leaflet.js + OSM", "Interactive GPS Mapping")
    ]
    for i, (title, desc) in enumerate(techs):
        x = Inches(0.8 + i * 2.983)
        add_content_card(slide, x, Inches(1.7), Inches(2.8), Inches(1.25))
        box = slide.shapes.add_textbox(x + Inches(0.15), Inches(1.82), Inches(2.5), Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        r1.font.size = Pt(13)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        
        p2 = tf.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = GRAY_TEXT

    # Architecture diagram container
    add_content_card(slide, Inches(0.8), Inches(3.15), Inches(11.733), Inches(3.75))
    
    # Layer 1
    add_content_card(slide, Inches(1.1), Inches(3.4), Inches(5.4), Inches(1.2), bg_color=LIGHT_BLUE_BG, border_color=BLUE_BORDER)
    b1 = slide.shapes.add_textbox(Inches(1.25), Inches(3.5), Inches(5.1), Inches(1.0))
    tf1 = b1.text_frame
    p = tf1.paragraphs[0]
    r = p.add_run()
    r.text = "Client Layer (Frontend UI)"
    r.font.bold = True
    r.font.size = Pt(12)
    r.font.color.rgb = BLUE
    p_d = tf1.add_paragraph()
    r_d = p_d.add_run()
    r_d.text = "Vanilla JS, Leaflet.js Mapping, Chart.js, HTML5 Geolocation API"
    r_d.font.size = Pt(10)
    r_d.font.color.rgb = GRAY_TEXT

    add_content_card(slide, Inches(6.833), Inches(3.4), Inches(5.4), Inches(1.2), bg_color=ORANGE_BG, border_color=ORANGE_BORDER)
    b2 = slide.shapes.add_textbox(Inches(6.983), Inches(3.5), Inches(5.1), Inches(1.0))
    tf2 = b2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Real-Time Layer (Firebase Web SDK)"
    r2.font.bold = True
    r2.font.size = Pt(12)
    r2.font.color.rgb = ORANGE_TEXT
    p2_d = tf2.add_paragraph()
    r2_d = p2_d.add_run()
    r2_d.text = "Firestore Snapshot Listeners for real-time GPS location pan & live bus group chat"
    r2_d.font.size = Pt(10)
    r2_d.font.color.rgb = GRAY_TEXT

    # Arrow divider
    arr_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(11.733), Inches(0.35))
    tf_arr = arr_box.text_frame
    p_arr = tf_arr.paragraphs[0]
    p_arr.alignment = PP_ALIGN.CENTER
    r_arr = p_arr.add_run()
    r_arr.text = "⬇️ REST API Queries & Real-Time Data Synchronization ⬇️"
    r_arr.font.size = Pt(11)
    r_arr.font.bold = True
    r_arr.font.color.rgb = BLUE

    # Layer 2
    add_content_card(slide, Inches(1.1), Inches(5.2), Inches(5.4), Inches(1.4), bg_color=GREEN_BG, border_color=GREEN_BORDER)
    b3 = slide.shapes.add_textbox(Inches(1.25), Inches(5.3), Inches(5.1), Inches(1.2))
    tf3 = b3.text_frame
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Backend API Layer (ASP.NET Core 8 C#)"
    r3.font.bold = True
    r3.font.size = Pt(12)
    r3.font.color.rgb = GREEN_TEXT
    p3_d = tf3.add_paragraph()
    r3_d = p3_d.add_run()
    r3_d.text = "REST Controllers, JWT Bearer Auth tokens, Role Authorization, Haversine Distance Engine"
    r3_d.font.size = Pt(10)
    r3_d.font.color.rgb = GRAY_TEXT

    add_content_card(slide, Inches(6.833), Inches(5.2), Inches(5.4), Inches(1.4), bg_color=PURPLE_BG, border_color=PURPLE_BORDER)
    b4 = slide.shapes.add_textbox(Inches(6.983), Inches(5.3), Inches(5.1), Inches(1.2))
    tf4 = b4.text_frame
    p4 = tf4.paragraphs[0]
    r4 = p4.add_run()
    r4.text = "Data & Third-Party Services"
    r4.font.bold = True
    r4.font.size = Pt(12)
    r4.font.color.rgb = DARK_TEXT
    p4_d = tf4.add_paragraph()
    r4_d = p4_d.add_run()
    r4_d.text = "Cloud Firestore NoSQL database + Brevo SMTP API for 6-digit OTP verification emails"
    r4_d.font.size = Pt(10)
    r4_d.font.color.rgb = GRAY_TEXT

    # ==================== SLIDE 4: DEVELOPMENT METHODOLOGY & WORKFLOW ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "Development Methodology & Workflow", "", "04")
    add_description(slide, "Iterative Agile development process emphasizing high reliability, real-time synchronization, and modularity.")

    # Phase 1-3
    add_content_card(slide, Inches(0.8), Inches(1.7), Inches(5.666), Inches(5.2))
    ttl1 = slide.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.0), Inches(0.4))
    r1 = ttl1.text_frame.paragraphs[0].add_run()
    r1.text = "📋 Agile Phases (1 — 3)"
    r1.font.size = Pt(17)
    r1.font.bold = True
    r1.font.color.rgb = DARK_TEXT

    p1_items = [
        ("Phase 1: Architecture & NoSQL Schema:", "Designed normalized User, Bus, Route, and Notice models in Firestore with indexed queries."),
        ("Phase 2: Responsive Prototyping:", "Built mobile-first interfaces with full multi-theme support (Classic Navy, Terracotta, Oceanic)."),
        ("Phase 3: Secure REST API:", "Implemented ASP.NET Core endpoints with Brevo SMTP 6-digit OTP verification and JWT claims.")
    ]
    add_rich_bullet_points(slide, p1_items, Inches(1.1), Inches(2.45), Inches(5.1), Inches(4.2), font_size=11, bullet_type="check")

    # Phase 4-6
    add_content_card(slide, Inches(6.866), Inches(1.7), Inches(5.666), Inches(5.2), bg_color=GREEN_BG, border_color=GREEN_BORDER)
    ttl2 = slide.shapes.add_textbox(Inches(7.166), Inches(1.9), Inches(5.0), Inches(0.4))
    r2 = ttl2.text_frame.paragraphs[0].add_run()
    r2.text = "🚀 Agile Phases (4 — 6)"
    r2.font.size = Pt(17)
    r2.font.bold = True
    r2.font.color.rgb = GREEN_TEXT

    p2_items = [
        ("Phase 4: GPS Tracking & ETA Calculation:", "Integrated HTML5 Geolocation API with Leaflet pan and Haversine distance engine."),
        ("Phase 5: Real-time Messaging & Broadcast:", "Built bus-wise community group chats and urgent notice board."),
        ("Phase 6: Verification & Testing:", "Executed xUnit backend unit tests and Playwright E2E integration test suites.")
    ]
    add_rich_bullet_points(slide, p2_items, Inches(7.166), Inches(2.45), Inches(5.1), Inches(4.2), font_size=11, bullet_type="check")

    # ==================== SLIDE 5: WALKTHROUGH: ONBOARDING & AUTH ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "Walkthrough: Onboarding & Multi-Role Authentication", "", "05")
    add_description(slide, "Role-based access system (Student, Driver, Admin) with secure registration, Brevo SMTP OTP verification, and JWT session handling.")

    add_screenshot_card(slide, "Screenshot 2026-08-16 110246.png", "Role Selection — Student / Driver / Admin", Inches(0.8), Inches(1.7), Inches(5.666), Inches(4.8))
    add_screenshot_card(slide, "Screenshot 2026-08-16 110301.png", "Secure Multi-Role Login Screen", Inches(6.866), Inches(1.7), Inches(5.666), Inches(4.8))

    # ==================== SLIDE 6: WALKTHROUGH: GPS & ETA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "Walkthrough: Real-Time GPS Tracking & Dynamic ETA", "", "06")
    add_description(slide, "Live Leaflet map pins bus coordinates updated every 5s. Haversine engine calculates exact stop distances and arrival estimates.")

    add_screenshot_card(slide, "Screenshot 2026-08-16 105418.png", "Live Map Tracking: Real-Time Location & Stop ETA Panel", Inches(0.8), Inches(1.7), Inches(11.733), Inches(4.8))

    # ==================== SLIDE 7: WALKTHROUGH: DRIVER CONSOLE & CHAT ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "Walkthrough: Driver Console & Community Chat", "", "07")
    add_description(slide, "Driver dashboard with single-tap trip start, delay broadcast, and bilingual UI, paired with real-time bus-wise group chat.")

    add_screenshot_card(slide, "WhatsApp Image 2026-08-16 at 10.55.43 AM (1).jpeg", "Driver Trip Controller & Delay Alerts (Mobile)", Inches(0.8), Inches(1.7), Inches(5.666), Inches(4.8))
    add_screenshot_card(slide, "Screenshot 2026-08-16 110042.png", "Bus-Wise Real-Time Passenger Community Chat", Inches(6.866), Inches(1.7), Inches(5.666), Inches(4.8))

    # ==================== SLIDE 8: WALKTHROUGH: MAP ROUTE BUILDER & FLEET CRUD ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "Walkthrough: Map Route Builder & Fleet Management", "", "08")
    add_description(slide, "Admin map coordinator enables clicking and dragging stop coordinates directly on OSM, alongside complete bus fleet CRUD operations.")

    add_screenshot_card(slide, "Screenshot 2026-08-16 103443.png", "Interactive Map Route Coordinator & Stop Editor", Inches(0.8), Inches(1.7), Inches(5.666), Inches(4.8))
    add_screenshot_card(slide, "Screenshot 2026-08-16 102304.png", "Bus Fleet Add / Modify / Delete CRUD Console", Inches(6.866), Inches(1.7), Inches(5.666), Inches(4.8))

    # ==================== SLIDE 9: WALKTHROUGH: ADMIN ANALYTICS & SAFETY ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "Walkthrough: Admin Analytics & Safety Management", "", "09")
    add_description(slide, "Admins monitor user statistics with Chart.js, broadcast urgent notices, manage user roles, and resolve student safety reports.")

    add_screenshot_card(slide, "Screenshot 2026-08-16 103910.png", "Chart.js Visual Analytics Dashboard", Inches(0.8), Inches(1.7), Inches(5.666), Inches(4.8))
    add_screenshot_card(slide, "Screenshot 2026-08-16 104456.png", "Student Safety & Bug Reporting Pipeline", Inches(6.866), Inches(1.7), Inches(5.666), Inches(4.8))

    # ==================== SLIDE 10: DIFFERENTIATORS & IMPACT ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "Key Differentiators & Measurable Impact", "", "10")
    add_description(slide, "SmartBus combines mathematical precision with seamless real-time interaction to deliver a superior campus commute.")

    # Differentiators
    add_content_card(slide, Inches(0.8), Inches(1.7), Inches(5.666), Inches(5.2), bg_color=GREEN_BG, border_color=GREEN_BORDER)
    ttl1 = slide.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.0), Inches(0.4))
    r1 = ttl1.text_frame.paragraphs[0].add_run()
    r1.text = "✨ Key Differentiators"
    r1.font.size = Pt(17)
    r1.font.bold = True
    r1.font.color.rgb = GREEN_TEXT

    diff_items = [
        ("🧮", "Haversine Formula ETA:", "Mathematical calculation based on geographical GPS coordinates."),
        ("💬", "Bus-Wise Chatrooms:", "Real-time Firestore communication per active bus."),
        ("🌐", "Bilingual Driver UI:", "Full English & বাংলা localization for seamless driver usability."),
        ("🗺️", "Visual Route Builder:", "Click-to-place stops directly on Leaflet OpenStreetMap.")
    ]
    add_rich_bullet_points(slide, diff_items, Inches(1.1), Inches(2.45), Inches(5.1), Inches(4.2), font_size=11, bullet_type="icon")

    # Impact
    add_content_card(slide, Inches(6.866), Inches(1.7), Inches(5.666), Inches(5.2))
    ttl2 = slide.shapes.add_textbox(Inches(7.166), Inches(1.9), Inches(5.0), Inches(0.4))
    r2 = ttl2.text_frame.paragraphs[0].add_run()
    r2.text = "📈 Measurable Impact"
    r2.font.size = Pt(17)
    r2.font.bold = True
    r2.font.color.rgb = DARK_TEXT

    impact_items = [
        ("⚡", "Zero Blind Waiting:", "Students track bus locations before heading to stops, eliminating idle wait."),
        ("📢", "Instant Delay Alerts:", "Push notices reduce campus confusion during unexpected traffic."),
        ("🛡️", "Enhanced Student Safety:", "Direct reporting pipeline with immediate admin response."),
        ("📊", "Data-Driven Transit:", "Analytics help optimize bus schedules and fleet allocation.")
    ]
    add_rich_bullet_points(slide, impact_items, Inches(7.166), Inches(2.45), Inches(5.1), Inches(4.2), font_size=11, bullet_type="icon")

    # ==================== SLIDE 11: BUSINESS MODEL & MONETIZATION ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "Business Model & Monetization Strategy", "", "11")
    add_description(slide, "A high-margin B2B SaaS transit management platform converting university fleets into recurring subscription revenue.")

    # 4 Cards layout
    card_w = Inches(2.78)
    gap = Inches(0.20)
    top_pos = Inches(1.7)
    card_h = Inches(5.2)

    # 1. Target Market
    x1 = Inches(0.8)
    add_content_card(slide, x1, top_pos, card_w, card_h)
    t1 = slide.shapes.add_textbox(x1 + Inches(0.2), Inches(1.9), card_w - Inches(0.4), Inches(0.4))
    r1 = t1.text_frame.paragraphs[0].add_run()
    r1.text = "🎯 Target Market"
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = DARK_TEXT

    mkt_items = [
        ("150+ Universities:", "Public & Private campus networks (BUBT, NSU, BRAC, DU)."),
        ("K-12 Schools:", "Student safety & parental live tracking subscriptions."),
        ("Corporate Fleets:", "IT & industrial employee shift shuttle fleets.")
    ]
    add_rich_bullet_points(slide, mkt_items, x1 + Inches(0.2), Inches(2.45), card_w - Inches(0.4), Inches(4.2), font_size=10, bullet_type="check")

    # 2. Value Proposition
    x2 = x1 + card_w + gap
    add_content_card(slide, x2, top_pos, card_w, card_h, bg_color=GREEN_BG, border_color=GREEN_BORDER)
    t2 = slide.shapes.add_textbox(x2 + Inches(0.2), Inches(1.9), card_w - Inches(0.4), Inches(0.4))
    r2 = t2.text_frame.paragraphs[0].add_run()
    r2.text = "💎 Value Proposition"
    r2.font.size = Pt(15)
    r2.font.bold = True
    r2.font.color.rgb = GREEN_TEXT

    val_items = [
        ("Zero Hardware Cost:", "100% Smartphone GPS based — no GPS tracker box needed."),
        ("Zero Blind Waiting:", "Real-time ETA & live delay alerts eliminate idle wait time."),
        ("20-30% Fleet Savings:", "Route & schedule analytics optimize fuel & transit costs.")
    ]
    add_rich_bullet_points(slide, val_items, x2 + Inches(0.2), Inches(2.45), card_w - Inches(0.4), Inches(4.2), font_size=10, bullet_type="check")

    # 3. Revenue Streams
    x3 = x2 + card_w + gap
    add_content_card(slide, x3, top_pos, card_w, card_h, bg_color=ORANGE_BG, border_color=ORANGE_BORDER)
    t3 = slide.shapes.add_textbox(x3 + Inches(0.2), Inches(1.9), card_w - Inches(0.4), Inches(0.4))
    r3 = t3.text_frame.paragraphs[0].add_run()
    r3.text = "💰 Revenue Streams"
    r3.font.size = Pt(15)
    r3.font.bold = True
    r3.font.color.rgb = ORANGE_TEXT

    rev_items = [
        ("B2B SaaS Subscription:", "$15 - $25 (৳1,500 - ৳2,500) / bus / month recurring fee."),
        ("White-Label Solution:", "৳50,000 custom campus setup & branded portal license."),
        ("Hyperlocal Campus Ads:", "Targeted deals for local food courts, bookshops & cafes.")
    ]
    add_rich_bullet_points(slide, rev_items, x3 + Inches(0.2), Inches(2.45), card_w - Inches(0.4), Inches(4.2), font_size=10, bullet_type="check")

    # 4. Unit Economics & Scale
    x4 = x3 + card_w + gap
    add_content_card(slide, x4, top_pos, card_w, card_h)
    t4 = slide.shapes.add_textbox(x4 + Inches(0.2), Inches(1.9), card_w - Inches(0.4), Inches(0.4))
    r4 = t4.text_frame.paragraphs[0].add_run()
    r4.text = "📈 Unit Economics"
    r4.font.size = Pt(15)
    r4.font.bold = True
    r4.font.color.rgb = DARK_TEXT

    scl_items = [
        ("85%+ Gross Margin:", "Zero map licensing fee (Leaflet/OSM) & low server cost."),
        ("Scalable ARR:", "10 Universities (~150 buses) = ৳36 Lakhs / Year recurring."),
        ("Rapid Onboarding:", "Full university campus fleet goes live in < 10 minutes.")
    ]
    add_rich_bullet_points(slide, scl_items, x4 + Inches(0.2), Inches(2.45), card_w - Inches(0.4), Inches(4.2), font_size=10, bullet_type="check")

    # ==================== SLIDE 12: CONCLUSION & FUTURE SCOPE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "Conclusion & Future Roadmap", "", "12")
    add_description(slide, "A complete, tested, and reliable university transit solution with clear paths for future expansion.")

    # Summary
    add_content_card(slide, Inches(0.8), Inches(1.7), Inches(5.666), Inches(5.2))
    ttl1 = slide.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.0), Inches(0.4))
    r1 = ttl1.text_frame.paragraphs[0].add_run()
    r1.text = "🏁 Project Summary"
    r1.font.size = Pt(17)
    r1.font.bold = True
    r1.font.color.rgb = DARK_TEXT

    sum_items = [
        ("Successfully developed and integrated a real-time GPS transit platform for university campuses."),
        ("Integrated ASP.NET Core REST API, Firestore real-time sync, and Brevo OTP authentication."),
        ("Validated through comprehensive xUnit backend testing and Playwright E2E browser automation.")
    ]
    add_rich_bullet_points(slide, sum_items, Inches(1.1), Inches(2.55), Inches(5.1), Inches(4.0), font_size=11, bullet_type="check")

    # Roadmap
    add_content_card(slide, Inches(6.866), Inches(1.7), Inches(5.666), Inches(5.2), bg_color=GREEN_BG, border_color=GREEN_BORDER)
    ttl2 = slide.shapes.add_textbox(Inches(7.166), Inches(1.9), Inches(5.0), Inches(0.4))
    r2 = ttl2.text_frame.paragraphs[0].add_run()
    r2.text = "🔮 Future Roadmap"
    r2.font.size = Pt(17)
    r2.font.bold = True
    r2.font.color.rgb = GREEN_TEXT

    fut_items = [
        ("IoT Passenger Crowding Sensor:", "Real-time bus seat occupancy tracking via smart sensors."),
        ("Automated Speed Alerts:", "Instant notifications for bus overspeeding events to ensure safety."),
        ("Offline SMS Fallback:", "Check bus timings and alerts via standard offline SMS queries."),
        ("Digital Ticketing Engine:", "Integrated payment gateway for special inter-campus routes.")
    ]
    add_rich_bullet_points(slide, fut_items, Inches(7.166), Inches(2.45), Inches(5.1), Inches(4.2), font_size=11, bullet_type="check")

    # ==================== SLIDE 13: THANK YOU ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)

    # Thank you header
    ty_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.4), Inches(10.333), Inches(1.0))
    tf = ty_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You!"
    run.font.size = Pt(50)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Outfit"

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.333), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Smart Campus Bus Tracking System"
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(200, 220, 255)
    run.font.name = "Inter"

    # Project Info
    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.3), Inches(10.333), Inches(2.4))
    tf = info_box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "Software Development Project 400"
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Intake 51 | Section 3"
    r2.font.size = Pt(13)
    r2.font.color.rgb = RGBColor(190, 215, 250)

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_after = Pt(8)
    r3 = p3.add_run()
    r3.text = "Supervised by "
    r3.font.size = Pt(14)
    r3.font.color.rgb = RGBColor(200, 220, 255)
    r3_s = p3.add_run()
    r3_s.text = "Humayra Ahmed"
    r3_s.font.bold = True
    r3_s.font.size = Pt(14)
    r3_s.font.color.rgb = WHITE
    r3_t = p3.add_run()
    r3_t.text = ", Assistant Professor"
    r3_t.font.size = Pt(14)
    r3_t.font.color.rgb = RGBColor(200, 220, 255)

    p4 = tf.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    p4.space_after = Pt(14)
    r4 = p4.add_run()
    r4.text = "Md Arif Khan  •  Karnia Binte Rafique  •  Suraiya Karim  •  Prosenjit Biswas  •  Proshanta Saha"
    r4.font.size = Pt(13)
    r4.font.color.rgb = RGBColor(220, 235, 255)

    p5 = tf.add_paragraph()
    p5.alignment = PP_ALIGN.CENTER
    r5 = p5.add_run()
    r5.text = "Questions & Feedback Welcome 🙌"
    r5.font.size = Pt(17)
    r5.font.bold = True
    r5.font.color.rgb = RGBColor(147, 197, 253)

    # Save
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Smart_Campus_Bus_Tracking_System.pptx')
    prs.save(output_path)
    print(f"[OK] Presentation saved to: {output_path}")
    print(f"[INFO] Total slides: {len(prs.slides)}")

if __name__ == '__main__':
    create_presentation()
